# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx>=1.0,<2", "click>=8.3,<9"]
# ///
"""PowerPoint (.pptx) operations tool.

Commands: create, info, lint, render, extract.

Usage:
    uv run pptx_tool.py COMMAND [OPTIONS]
"""

import json
import re
import sys
from html import escape
from math import ceil
from pathlib import Path

import click
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def write_output(text: str, output_path: str | None) -> None:
    """Write text to file or stdout."""
    if output_path:
        Path(output_path).write_text(text, encoding="utf-8")
        click.echo(f"Output written to {output_path}", err=True)
    else:
        click.echo(text)


def _load_json_input(value: str) -> object:
    candidate = Path(value)
    if candidate.exists() and candidate.is_file():
        return json.loads(candidate.read_text(encoding="utf-8"))
    return json.loads(value)


def _rgb(value: object, default: str = "000000") -> RGBColor:
    text = str(value if value is not None else default).strip().lstrip("#")
    if len(text) != 6:
        raise ValueError(f"invalid RGB color: {value}")
    return RGBColor.from_string(text.upper())


def _inches(spec: dict[str, object], key: str, default: float) -> int:
    return Inches(float(spec.get(key, default)))


def _layout_for(prs: Presentation, value: object | None, fallback: int):
    if value is not None:
        if isinstance(value, int) or (isinstance(value, str) and value.isdigit()):
            index = int(value)
            if not 0 <= index < len(prs.slide_layouts):
                raise ValueError(f"slide layout index out of range: {index}")
            return prs.slide_layouts[index]
        wanted = str(value).casefold()
        for layout in prs.slide_layouts:
            if layout.name.casefold() == wanted:
                return layout
        raise ValueError(f"slide layout not found: {value}")
    if not prs.slide_layouts:
        raise ValueError("template has no slide layouts")
    return prs.slide_layouts[min(fallback, len(prs.slide_layouts) - 1)]


def _apply_text_frame(tf, spec: dict[str, object], theme: dict[str, object]) -> None:
    tf.word_wrap = bool(spec.get("wrap", True))
    vertical = str(spec.get("vertical", "top"))
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(vertical, MSO_ANCHOR.TOP)
    text = str(spec.get("text", ""))
    lines = text.split("\n") or [""]
    tf.clear()
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }.get(str(spec.get("alignment", "left")), PP_ALIGN.LEFT)
        paragraph.level = int(spec.get("level", 0))
        for run in paragraph.runs:
            run.font.name = str(spec.get("fontName", theme.get("fontName", "Arial")))
            run.font.size = Pt(float(spec.get("fontSize", theme.get("fontSize", 20))))
            run.font.bold = bool(spec.get("bold", False))
            run.font.italic = bool(spec.get("italic", False))
            run.font.color.rgb = _rgb(spec.get("color", theme.get("textColor", "1F2937")))


def _shape_type(value: object) -> MSO_SHAPE:
    mapping = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "roundedRectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "chevron": MSO_SHAPE.CHEVRON,
        "diamond": MSO_SHAPE.DIAMOND,
        "hexagon": MSO_SHAPE.HEXAGON,
    }
    if str(value) not in mapping:
        raise ValueError(f"unsupported shape type: {value}")
    return mapping[str(value)]


def _add_element(slide, spec: dict[str, object], theme: dict[str, object]) -> None:
    element_type = str(spec.get("type", "text"))
    x, y = _inches(spec, "x", 1), _inches(spec, "y", 1)
    width, height = _inches(spec, "width", 4), _inches(spec, "height", 1)
    if element_type == "text":
        shape = slide.shapes.add_textbox(x, y, width, height)
        _apply_text_frame(shape.text_frame, spec, theme)
    elif element_type == "shape":
        shape = slide.shapes.add_shape(_shape_type(spec.get("shape", "rectangle")), x, y, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(spec.get("fill", theme.get("accentColor", "4F46E5")))
        shape.line.color.rgb = _rgb(spec.get("lineColor", spec.get("fill", theme.get("accentColor", "4F46E5"))))
        if spec.get("text") is not None:
            _apply_text_frame(shape.text_frame, spec, theme)
    elif element_type == "image":
        if not spec.get("path"):
            raise ValueError("image element requires path")
        slide.shapes.add_picture(str(spec["path"]), x, y, width=width, height=height)
        return
    elif element_type == "table":
        rows = spec.get("rows")
        if not isinstance(rows, list) or not rows or any(not isinstance(row, list) for row in rows):
            raise ValueError("table rows must be a non-empty two-dimensional array")
        column_count = max(len(row) for row in rows)
        table = slide.shapes.add_table(len(rows), column_count, x, y, width, height).table
        for row_index, row in enumerate(rows):
            for column_index, value in enumerate(row):
                cell = table.cell(row_index, column_index)
                cell.text = str(value.get("text", "") if isinstance(value, dict) else value)
                cell.fill.solid()
                is_header = row_index == 0 and bool(spec.get("header", True))
                cell.fill.fore_color.rgb = _rgb(
                    value.get("fill") if isinstance(value, dict) and value.get("fill") else
                    (spec.get("headerFill", theme.get("accentColor", "4F46E5")) if is_header else spec.get("fill", "FFFFFF"))
                )
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(float(spec.get("fontSize", 14)))
                        run.font.bold = is_header
                        run.font.color.rgb = _rgb("FFFFFF" if is_header else theme.get("textColor", "1F2937"))
        return
    elif element_type == "chart":
        categories = spec.get("categories")
        series = spec.get("series")
        if not isinstance(categories, list) or not isinstance(series, list) or not series:
            raise ValueError("chart requires categories and a non-empty series array")
        data = CategoryChartData()
        data.categories = [str(value) for value in categories]
        for item in series:
            if not isinstance(item, dict) or not isinstance(item.get("values"), list):
                raise ValueError("chart series require name and values")
            data.add_series(str(item.get("name", "Series")), item["values"])
        chart_type = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }.get(str(spec.get("chartType", "column")))
        if chart_type is None:
            raise ValueError(f"unsupported chart type: {spec.get('chartType')}")
        chart = slide.shapes.add_chart(chart_type, x, y, width, height, data).chart
        chart.has_title = bool(spec.get("title"))
        if chart.has_title:
            chart.chart_title.text_frame.text = str(spec["title"])
        chart.has_legend = bool(spec.get("legend", True))
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        return
    else:
        raise ValueError(f"unsupported slide element type: {element_type}")
    if spec.get("rotation") is not None:
        shape.rotation = float(spec["rotation"])


def _shape_bounds(shape) -> dict[str, object]:
    return {
        "name": shape.name,
        "type": str(shape.shape_type),
        "x": round(shape.left / 914400, 3),
        "y": round(shape.top / 914400, 3),
        "width": round(shape.width / 914400, 3),
        "height": round(shape.height / 914400, 3),
        "text": shape.text[:200] if getattr(shape, "has_text_frame", False) else None,
        "has_table": bool(getattr(shape, "has_table", False)),
        "has_chart": bool(getattr(shape, "has_chart", False)),
    }


def _lint_presentation(prs: Presentation) -> list[dict[str, object]]:
    issues: list[dict[str, object]] = []
    slide_width, slide_height = prs.slide_width / 914400, prs.slide_height / 914400
    for slide_index, slide in enumerate(prs.slides, start=1):
        shapes = [_shape_bounds(shape) for shape in slide.shapes]
        for shape in shapes:
            if shape["x"] < 0 or shape["y"] < 0 or shape["x"] + shape["width"] > slide_width or shape["y"] + shape["height"] > slide_height:
                issues.append({"slide": slide_index, "type": "off-page", "shape": shape["name"], "bounds": shape})
            text = shape.get("text")
            if text:
                estimated_lines = max(1, ceil(len(str(text)) / max(1, int(float(shape["width"]) * 9))))
                if estimated_lines * 0.3 > float(shape["height"]):
                    issues.append({"slide": slide_index, "type": "text-overflow", "shape": shape["name"], "estimatedLines": estimated_lines})
        non_placeholders = [shape for shape, raw in zip(shapes, slide.shapes) if not raw.is_placeholder]
        for first_index, first in enumerate(non_placeholders):
            for second in non_placeholders[first_index + 1:]:
                overlap_width = min(first["x"] + first["width"], second["x"] + second["width"]) - max(first["x"], second["x"])
                overlap_height = min(first["y"] + first["height"], second["y"] + second["height"]) - max(first["y"], second["y"])
                if overlap_width > 0.03 and overlap_height > 0.03:
                    issues.append({"slide": slide_index, "type": "overlap", "shapes": [first["name"], second["name"]], "area": round(overlap_width * overlap_height, 3)})
    return issues


@click.group()
def cli() -> None:
    """PowerPoint (.pptx) operations tool."""
    pass


@cli.command()
@click.option("--from-file", type=click.Path(exists=True, dir_okay=False), default=None, help="Input markdown or JSON file.")
@click.option("--text", type=str, default=None, help="Inline markdown/text content for slides.")
@click.option("--json-data", type=str, default=None, help="JSON string defining slides structure.")
@click.option("--title", type=str, default=None, help="Presentation title (creates title slide).")
@click.option("--template", type=click.Path(exists=True, dir_okay=False), default=None, help="Template .pptx file to use.")
@click.option("-o", "--output", type=click.Path(), required=True, help="Output .pptx file path.")
def create(from_file: str | None, text: str | None, json_data: str | None, title: str | None, template: str | None, output: str) -> None:
    """Create a PowerPoint presentation from markdown, text, or JSON.

    Markdown format: Use '---' to separate slides. '#' for titles, body text below.

    JSON accepts a legacy slide array or a structured presentation object with
    title, subtitle, properties, slideSize, theme, and slides. Structured slides
    can contain text, shape, image, table, and chart elements.
    """
    if from_file is None and text is None and json_data is None and title is None:
        click.echo("Error: provide --from-file, --text, --json-data, or --title.", err=True)
        sys.exit(1)

    try:
        if template:
            prs = Presentation(template)
        else:
            prs = Presentation()

        slides_data: list[dict[str, object]] = []
        presentation_title = title
        presentation_subtitle = ""
        theme: dict[str, object] = {}
        properties: dict[str, object] = {}
        slide_size: dict[str, object] = {}

        if json_data:
            parsed = _load_json_input(json_data)
            if isinstance(parsed, list):
                slides_data = parsed
            elif isinstance(parsed, dict):
                raw_slides = parsed.get("slides", [])
                if not isinstance(raw_slides, list):
                    raise ValueError("JSON 'slides' must be an array")
                slides_data = raw_slides
                presentation_title = title or (str(parsed["title"]) if parsed.get("title") else None)
                presentation_subtitle = str(parsed.get("subtitle", ""))
                raw_theme = parsed.get("theme", {})
                raw_properties = parsed.get("properties", {})
                raw_slide_size = parsed.get("slideSize", {})
                if not isinstance(raw_theme, dict) or not isinstance(raw_properties, dict) or not isinstance(raw_slide_size, dict):
                    raise ValueError("theme, properties, and slideSize must be objects")
                theme = raw_theme
                properties = raw_properties
                slide_size = raw_slide_size
            else:
                raise ValueError("JSON must be a slide array or presentation object")

        elif from_file:
            content = Path(from_file).read_text(encoding="utf-8")
            if from_file.endswith(".json"):
                parsed = json.loads(content)
                if isinstance(parsed, list):
                    slides_data = parsed
                elif isinstance(parsed, dict):
                    raw_slides = parsed.get("slides", [])
                    if not isinstance(raw_slides, list):
                        raise ValueError("JSON 'slides' must be an array")
                    slides_data = raw_slides
                    presentation_title = title or (str(parsed["title"]) if parsed.get("title") else None)
                    presentation_subtitle = str(parsed.get("subtitle", ""))
                    theme = parsed.get("theme", {})
                    properties = parsed.get("properties", {})
                    slide_size = parsed.get("slideSize", {})
                    if not isinstance(theme, dict) or not isinstance(properties, dict) or not isinstance(slide_size, dict):
                        raise ValueError("theme, properties, and slideSize must be objects")
                else:
                    raise ValueError("JSON must be a slide array or presentation object")
            else:
                slides_data = _parse_markdown_slides(content)

        elif text:
            slides_data = _parse_markdown_slides(text)

        if slide_size:
            prs.slide_width = Inches(float(slide_size.get("width", prs.slide_width / 914400)))
            prs.slide_height = Inches(float(slide_size.get("height", prs.slide_height / 914400)))

        for key in ("title", "subject", "author", "keywords", "comments", "category"):
            if properties.get(key) is not None and hasattr(prs.core_properties, key):
                setattr(prs.core_properties, key, str(properties[key]))

        if presentation_title:
            _add_title_slide(prs, presentation_title, presentation_subtitle, theme)

        for slide_info in slides_data:
            if not isinstance(slide_info, dict):
                raise ValueError("each slide must be an object")
            _add_content_slide(prs, slide_info, theme)

        prs.save(output)
        click.echo(f"Presentation saved to {output} ({len(slides_data) + (1 if presentation_title else 0)} slides)", err=True)
    except json.JSONDecodeError as e:
        click.echo(f"Error parsing JSON: {e}", err=True)
        sys.exit(1)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


def _parse_markdown_slides(md_text: str) -> list[dict[str, str]]:
    """Parse markdown into slide data. Slides separated by '---'."""
    raw_slides = re.split(r"\n---\n|\n---$|^---\n", md_text)
    slides: list[dict[str, str]] = []

    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue

        lines = raw.split("\n")
        slide_title = ""
        body_lines: list[str] = []

        for line in lines:
            heading_match = re.match(r"^#{1,3}\s+(.+)$", line.strip())
            if heading_match and not slide_title:
                slide_title = heading_match.group(1)
            else:
                body_lines.append(line)

        body = "\n".join(body_lines).strip()
        slide: dict[str, str] = {}
        if slide_title:
            slide["title"] = slide_title
        if body:
            slide["body"] = body
        if slide:
            slides.append(slide)

    return slides


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "", theme: dict[str, object] | None = None) -> None:
    """Add a title slide."""
    theme = theme or {}
    layout = _layout_for(prs, theme.get("titleLayout"), 0)
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(theme.get("backgroundColor", "FFFFFF"))

    if slide.placeholders:
        if 0 in slide.placeholders:
            _apply_text_frame(slide.placeholders[0].text_frame, {
                "text": title,
                "fontName": theme.get("fontName", "Arial"),
                "fontSize": theme.get("titleFontSize", 32),
                "color": theme.get("titleColor", theme.get("textColor", "1F2937")),
                "bold": True,
            }, theme)
        if 1 in slide.placeholders and subtitle:
            _apply_text_frame(slide.placeholders[1].text_frame, {
                "text": subtitle,
                "fontName": theme.get("fontName", "Arial"),
                "fontSize": theme.get("subtitleFontSize", 18),
                "color": theme.get("mutedColor", "64748B"),
            }, theme)


def _add_content_slide(prs: Presentation, slide_info: dict[str, object], theme: dict[str, object]) -> None:
    """Add a content slide from slide info dict."""
    title = str(slide_info.get("title", ""))
    body = str(slide_info.get("body", ""))
    notes = str(slide_info.get("notes", ""))
    elements = slide_info.get("elements", [])
    if not isinstance(elements, list):
        raise ValueError("slide elements must be an array")

    if elements:
        fallback_layout = 5 if title else 6
    elif title and body:
        fallback_layout = 1
    elif title:
        fallback_layout = 5
    else:
        fallback_layout = 6
    layout = _layout_for(prs, slide_info.get("layout"), fallback_layout)

    slide = prs.slides.add_slide(layout)
    background = slide_info.get("background", theme.get("backgroundColor"))
    if background:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(background)

    if title and slide.placeholders and 0 in slide.placeholders:
        _apply_text_frame(slide.placeholders[0].text_frame, {
            "text": title,
            "fontName": theme.get("fontName", "Arial"),
            "fontSize": slide_info.get("titleFontSize", theme.get("titleFontSize", 28)),
            "color": slide_info.get("titleColor", theme.get("titleColor", theme.get("textColor", "1F2937"))),
            "bold": True,
        }, theme)

    if body:
        if slide.placeholders and 1 in slide.placeholders:
            tf = slide.placeholders[1].text_frame
        else:
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            tf = text_box.text_frame
            tf.word_wrap = True
        tf.clear()

        body_lines = body.split("\n")
        first_line = True
        for line in body_lines:
            stripped = line.strip()
            if not stripped:
                continue

            if first_line:
                p = tf.paragraphs[0]
                first_line = False
            else:
                p = tf.add_paragraph()

            bullet_match = re.match(r"^[-*+]\s+(.+)$", stripped)
            num_match = re.match(r"^\d+[.)]\s+(.+)$", stripped)

            if bullet_match:
                p.text = bullet_match.group(1)
                p.level = 0
            elif num_match:
                p.text = num_match.group(1)
                p.level = 0
            else:
                indent_match = re.match(r"^(\s+)[-*+]\s+(.+)$", line)
                indent_num_match = re.match(r"^(\s+)\d+[.)]\s+(.+)$", line)
                if indent_match:
                    indent_level = len(indent_match.group(1)) // 2
                    p.text = indent_match.group(2)
                    p.level = min(indent_level, 4)
                elif indent_num_match:
                    indent_level = len(indent_num_match.group(1)) // 2
                    p.text = indent_num_match.group(2)
                    p.level = min(indent_level, 4)
                else:
                    p.text = stripped

            for run in p.runs:
                run.font.name = str(theme.get("fontName", "Arial"))
                run.font.size = Pt(float(theme.get("fontSize", 20)))
                run.font.color.rgb = _rgb(theme.get("textColor", "1F2937"))

    for element in elements:
        if not isinstance(element, dict):
            raise ValueError("each slide element must be an object")
        _add_element(slide, element, theme)

    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("-o", "--output", type=click.Path(), default=None, help="Write output to file.")
def info(file: str, output: str | None) -> None:
    """Show presentation metadata and slide information."""
    try:
        prs = Presentation(file)

        lint_issues = _lint_presentation(prs)
        slides_info = []
        for i, slide in enumerate(prs.slides):
            shapes = [_shape_bounds(shape) for shape in slide.shapes]
            slide_data: dict[str, object] = {
                "number": i + 1,
                "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
                "shapes": shapes,
            }

            # Get title
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text

            # Count shapes by type
            slide_data["shape_count"] = len(slide.shapes)
            slide_data["table_count"] = sum(1 for shape in slide.shapes if getattr(shape, "has_table", False))
            slide_data["chart_count"] = sum(1 for shape in slide.shapes if getattr(shape, "has_chart", False))
            slide_data["image_count"] = sum(
                1 for shape in slide.shapes
                if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)
            )
            slide_data["lint_issue_count"] = sum(1 for issue in lint_issues if issue["slide"] == i + 1)

            # Notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_data["has_notes"] = True

            slides_info.append(slide_data)

        # Slide dimensions
        width_in = prs.slide_width / 914400  # EMU to inches
        height_in = prs.slide_height / 914400

        info_dict: dict[str, object] = {
            "file": str(Path(file).resolve()),
            "slide_count": len(prs.slides),
            "slide_width_inches": round(width_in, 2),
            "slide_height_inches": round(height_in, 2),
            "title": prs.core_properties.title,
            "author": prs.core_properties.author,
            "lint_issue_count": len(lint_issues),
            "slides": slides_info,
        }

        # Available layouts
        layouts = [layout.name for layout in prs.slide_layouts]
        info_dict["available_layouts"] = layouts

        write_output(json.dumps(info_dict, indent=2, default=str), output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--fail-on-issues", is_flag=True, help="Exit non-zero when layout issues are found.")
@click.option("-o", "--output", type=click.Path(), default=None, help="Write JSON output to file.")
def lint(file: str, fail_on_issues: bool, output: str | None) -> None:
    """Check for off-page shapes, likely text overflow, and overlaps."""
    try:
        issues = _lint_presentation(Presentation(file))
        write_output(json.dumps({
            "file": str(Path(file).resolve()),
            "issue_count": len(issues),
            "issues": issues,
        }, indent=2), output)
        if fail_on_issues and issues:
            raise click.exceptions.Exit(2)
    except click.exceptions.Exit:
        raise
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("-o", "--output", type=click.Path(), required=True, help="Output SVG contact sheet.")
@click.option("--columns", type=click.IntRange(1, 8), default=3, show_default=True)
def render(file: str, output: str, columns: int) -> None:
    """Render a deterministic SVG contact sheet for review without Office."""
    try:
        prs = Presentation(file)
        slide_width = prs.slide_width / 914400
        slide_height = prs.slide_height / 914400
        thumb_width = 320.0
        thumb_height = thumb_width * slide_height / slide_width
        gutter, label_height = 24.0, 28.0
        rows = max(1, ceil(len(prs.slides) / columns))
        canvas_width = columns * thumb_width + (columns + 1) * gutter
        canvas_height = rows * (thumb_height + label_height) + (rows + 1) * gutter
        svg = [
            f'<svg xmlns="http://www.w3.org/2000/svg" width="{canvas_width:.0f}" height="{canvas_height:.0f}" viewBox="0 0 {canvas_width:.0f} {canvas_height:.0f}">',
            '<rect width="100%" height="100%" fill="#e5e7eb"/>',
            '<style>text{font-family:Arial,sans-serif}.label{font-size:13px;fill:#334155}.content{font-size:10px;fill:#111827}</style>',
        ]
        for index, slide in enumerate(prs.slides):
            column, row = index % columns, index // columns
            origin_x = gutter + column * (thumb_width + gutter)
            origin_y = gutter + row * (thumb_height + label_height + gutter)
            svg.append(f'<rect x="{origin_x:.2f}" y="{origin_y:.2f}" width="{thumb_width:.2f}" height="{thumb_height:.2f}" rx="3" fill="#ffffff" stroke="#94a3b8"/>')
            for shape in slide.shapes:
                bounds = _shape_bounds(shape)
                x = origin_x + float(bounds["x"]) / slide_width * thumb_width
                y = origin_y + float(bounds["y"]) / slide_height * thumb_height
                width = max(1.0, float(bounds["width"]) / slide_width * thumb_width)
                height = max(1.0, float(bounds["height"]) / slide_height * thumb_height)
                if getattr(shape, "has_chart", False):
                    fill, label = "#dbeafe", "[chart]"
                elif getattr(shape, "has_table", False):
                    fill, label = "#dcfce7", "[table]"
                elif shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
                    fill, label = "#fef3c7", "[image]"
                else:
                    fill = "#f8fafc"
                    label = str(bounds.get("text") or "")[:80].replace("\n", " ")
                svg.append(f'<rect x="{x:.2f}" y="{y:.2f}" width="{width:.2f}" height="{height:.2f}" fill="{fill}" fill-opacity="0.78" stroke="#cbd5e1" stroke-width="0.6"/>')
                if label and width > 16 and height > 8:
                    svg.append(f'<text class="content" x="{x + 3:.2f}" y="{y + min(12, height - 2):.2f}">{escape(label)}</text>')
            svg.append(f'<text class="label" x="{origin_x:.2f}" y="{origin_y + thumb_height + 19:.2f}">Slide {index + 1}</text>')
        svg.append("</svg>")
        Path(output).write_text("\n".join(svg), encoding="utf-8")
        click.echo(f"Contact sheet written to {output}", err=True)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


@cli.command()
@click.argument("file", type=click.Path(exists=True, dir_okay=False))
@click.option("--slide", type=int, default=None, help="Extract specific slide number (1-based).")
@click.option("--include-notes/--no-notes", default=True, help="Include speaker notes (default: yes).")
@click.option("-o", "--output", type=click.Path(), default=None, help="Write output to file.")
def extract(file: str, slide: int | None, include_notes: bool, output: str | None) -> None:
    """Extract text content from a PowerPoint presentation."""
    try:
        prs = Presentation(file)
        parts: list[str] = []

        slides_to_process = []
        if slide is not None:
            if 1 <= slide <= len(prs.slides):
                slides_to_process = [(slide - 1, prs.slides[slide - 1])]
            else:
                click.echo(f"Error: slide {slide} out of range (1-{len(prs.slides)}).", err=True)
                sys.exit(1)
        else:
            slides_to_process = list(enumerate(prs.slides))

        for i, s in slides_to_process:
            parts.append(f"--- Slide {i + 1} ---")

            # Title
            if s.shapes.title:
                parts.append(f"# {s.shapes.title.text}")

            # All text shapes
            for shape in s.shapes:
                if shape.has_text_frame:
                    # Skip title shape (already handled)
                    if shape == s.shapes.title:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            # Add indentation for bullet levels
                            indent = "  " * paragraph.level if paragraph.level else ""
                            parts.append(f"{indent}- {text}" if paragraph.level > 0 else text)

                # Table content
                if shape.has_table:
                    table = shape.table
                    parts.append("")
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        parts.append(" | ".join(cells))

            # Speaker notes
            if include_notes and s.has_notes_slide:
                notes_text = s.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"\n[Notes: {notes_text}]")

            parts.append("")

        write_output("\n".join(parts), output)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)


if __name__ == "__main__":
    cli()
